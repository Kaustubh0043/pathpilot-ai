import io
import PyPDF2
import docx
import pptx

def extract_text_from_bytes(file_bytes: bytes, filename: str) -> str:
    """Extracts raw text content from uploaded file bytes depending on file extension."""
    ext = filename.split(".")[-1].lower()
    text = ""

    try:
        if ext == "pdf":
            reader = PyPDF2.PdfReader(io.BytesIO(file_bytes))
            for page in reader.pages:
                page_text = page.extract_text()
                if page_text:
                    text += page_text + "\n"
        
        elif ext in ["docx", "doc"]:
            doc = docx.Document(io.BytesIO(file_bytes))
            for para in doc.paragraphs:
                if para.text:
                    text += para.text + "\n"
        
        elif ext in ["pptx", "ppt"]:
            prs = pptx.Presentation(io.BytesIO(file_bytes))
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        text += shape.text + "\n"
        
        elif ext in ["txt", "md"]:
            try:
                text = file_bytes.decode("utf-8")
            except UnicodeDecodeError:
                text = file_bytes.decode("latin-1")
        
        else:
            raise ValueError(f"Unsupported file extension: {ext}")
            
    except Exception as e:
        raise RuntimeError(f"Error parsing document {filename}: {str(e)}")

    return text.strip()
